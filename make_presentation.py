from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
FIFA_BLUE   = RGBColor(0x00, 0x41, 0x8A)   # deep FIFA blue
UEFA_RED    = RGBColor(0xC0, 0x00, 0x00)   # UEFA red
DARK_GREY   = RGBColor(0x2E, 0x2E, 0x2E)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT      = RGBColor(0x00, 0x70, 0xC0)
GREEN       = RGBColor(0x37, 0x86, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_GREY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """Blue top bar with title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=FIFA_BLUE)
    add_textbox(slide, title,
                Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.65),
                font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(0.72), Inches(12.6), Inches(0.35),
                    font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF), italic=True)


def add_bullet_box(slide, items, left, top, width, height,
                   title=None, title_color=FIFA_BLUE, bullet="▸",
                   font_size=15, bg_color=LIGHT_GREY, border=True):
    """A rounded-corner-style box with optional title and bullet list."""
    box = add_rect(slide, left, top, width, height,
                   fill_color=bg_color,
                   line_color=RGBColor(0xCC, 0xCC, 0xCC) if border else None)
    tx = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                   width - Inches(0.3), height - Inches(0.2))
    tf = tx.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_GREY
    return tx


def add_image(slide, path, left, top, width, height=None):
    if height:
        pic = slide.shapes.add_picture(path, left, top, width, height)
    else:
        pic = slide.shapes.add_picture(path, left, top, width)
    return pic


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

# Full background
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=FIFA_BLUE)
# White accent bar on the left
add_rect(s, 0, 0, Inches(0.5), SLIDE_H, fill_color=RGBColor(0xFF, 0xFF, 0xFF))
# Red accent strip
add_rect(s, Inches(0.5), 0, Inches(0.12), SLIDE_H, fill_color=UEFA_RED)

# Title
add_textbox(s,
    "UEFA vs FIFA Tie Break:\nWhich Creates More Suspense?",
    Inches(1.0), Inches(1.8), Inches(11.5), Inches(2.2),
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(s,
    "Alessandro Di Mattia  ·  Alex Krumer",
    Inches(1.0), Inches(4.1), Inches(11.5), Inches(0.5),
    font_size=20, color=RGBColor(0xCC, 0xDD, 0xFF))

add_textbox(s,
    "Supervisor Meeting  ·  June 2026",
    Inches(1.0), Inches(4.65), Inches(11.5), Inches(0.4),
    font_size=16, color=RGBColor(0xAA, 0xBB, 0xDD), italic=True)

# Bottom tag line
add_textbox(s,
    "Sports Economics  ·  Tournament Design  ·  Suspense Measurement",
    Inches(1.0), Inches(6.6), Inches(11.5), Inches(0.4),
    font_size=13, color=RGBColor(0x88, 0xAA, 0xCC))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Motivation
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Motivation", "Why tie-breaking rules matter for fans and revenues")

# Two columns
left_items = [
    "2022 FIFA World Cup: ~2.87 billion viewers; USD 6.3 bn in broadcasting & commercial rights",
    "Euro 2024: 5.7 billion worldwide TV viewers; 2.67 million stadium attendees across 51 matches",
    "Fan engagement translates directly into broadcasting fees, sponsorships, and ticket revenues",
]
right_items = [
    "Fans derive enjoyment from non-instrumental information: suspense & uncertainty (Ely, 2015; Zillmann, 2013)",
    "In-game dynamics — not just final outcome — drive demand (Alavy et al., 2010; Buraimo et al., 2020)",
    "Low-stakes games reduce interest; rules that preserve uncertainty extend attention across the tournament",
]

add_bullet_box(s, left_items,
               Inches(0.3), Inches(1.25), Inches(6.2), Inches(5.1),
               title="The Economic Stakes", bg_color=RGBColor(0xE8, 0xF0, 0xFE))
add_bullet_box(s, right_items,
               Inches(6.8), Inches(1.25), Inches(6.2), Inches(5.1),
               title="The Academic Insight", bg_color=RGBColor(0xFE, 0xEC, 0xE8))

add_textbox(s,
    "▶  Tie-breaking criteria are an under-studied institutional lever that shapes suspense for billions of fans.",
    Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.55),
    font_size=14, bold=True, color=FIFA_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Research Question & Contribution
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Research Question & Contribution")

# Central RQ box
add_rect(s, Inches(1.2), Inches(1.3), Inches(10.9), Inches(1.3),
         fill_color=FIFA_BLUE)
add_textbox(s,
    "Does the choice of tie-breaking criterion — FIFA goal difference vs UEFA head-to-head — affect the level of group-stage suspense in international football?",
    Inches(1.4), Inches(1.35), Inches(10.5), Inches(1.2),
    font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

contrib_items = [
    "New binary suspense indicator defined at the group level (not match level)",
    "Counterfactual design: apply both rules to the same historical data",
    "Empirical test using ~662 goals and ~10,800 minute-by-minute observations (WC 1986–2022; Euro 1984–2024)",
    "European Championship as a built-in placebo check",
    "Bridges tournament design & sports demand literatures (Csató, 2023; Buraimo et al., 2020)",
]
add_bullet_box(s, contrib_items,
               Inches(0.3), Inches(2.85), Inches(12.7), Inches(3.9),
               title="Our Contributions", font_size=15,
               bg_color=LIGHT_GREY)

add_textbox(s, "Main result: FIFA's goal-difference rule generates significantly more suspense in the World Cup (+4.2 pp, p < 0.01)",
            Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.45),
            font_size=14, bold=True, color=UEFA_RED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Data
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Data", "World Cup 1986–2022 and European Championship 1984–2024")

# WC box
add_rect(s, Inches(0.3), Inches(1.25), Inches(6.0), Inches(5.0),
         fill_color=RGBColor(0xE8, 0xF0, 0xFE),
         line_color=FIFA_BLUE, line_width=Pt(1.5))
add_textbox(s, "FIFA World Cup", Inches(0.5), Inches(1.35), Inches(5.6), Inches(0.5),
            font_size=18, bold=True, color=FIFA_BLUE)
wc_lines = [
    "Editions: 1986 – 2022  (10 tournaments)",
    "Groups: 6 (1986–1994) → 8 (1998–2022), 4 teams each",
    "Qualifying format: Top 2 (+ best 4 thirds in 6-group era)",
    "Goals scraped: 455   |   Minute rows: 6,789",
    "Points: 2 for win until 1992; 3 for win from 1994",
    "Elo ratings: from www.international-football.net",
    "Cross-checked with official FIFA sources",
]
for i, line in enumerate(wc_lines):
    add_textbox(s, f"▸  {line}", Inches(0.5), Inches(1.9) + Inches(0.48*i),
                Inches(5.6), Inches(0.46), font_size=14, color=DARK_GREY)

# Euro box
add_rect(s, Inches(7.0), Inches(1.25), Inches(6.0), Inches(5.0),
         fill_color=RGBColor(0xFE, 0xEC, 0xE8),
         line_color=UEFA_RED, line_width=Pt(1.5))
add_textbox(s, "UEFA European Championship", Inches(7.2), Inches(1.35), Inches(5.6), Inches(0.5),
            font_size=18, bold=True, color=UEFA_RED)
eu_lines = [
    "Editions: 1984 – 2024  (11 tournaments)",
    "Groups: 2 (1984–1992) → 4 (1996–2012) → 6 (2016–2024)",
    "Qualifying format: Top 2 (+ best 4 thirds in 6-group era)",
    "Goals scraped: 207   |   Minute rows: 4,044",
    "Points: same evolution as FIFA",
    "Elo ratings: same source",
    "Cross-checked with official UEFA sources",
]
for i, line in enumerate(eu_lines):
    add_textbox(s, f"▸  {line}", Inches(7.2), Inches(1.9) + Inches(0.48*i),
                Inches(5.6), Inches(0.46), font_size=14, color=DARK_GREY)

add_textbox(s, "Counterfactual: both tie-breaking rules applied to the same observed match data",
            Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.45),
            font_size=14, bold=True, color=FIFA_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Methodology: Suspense Indicator
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Methodology: Binary Suspense Indicator")

add_textbox(s,
    "Suspense = 1  if a single additional goal in any ongoing group match\nwould change which teams advance to the knockout stage",
    Inches(0.3), Inches(1.25), Inches(12.7), Inches(1.05),
    font_size=18, bold=True, color=FIFA_BLUE, align=PP_ALIGN.CENTER)

# Two tracking modes
add_rect(s, Inches(0.3), Inches(2.5), Inches(5.8), Inches(1.8),
         fill_color=LIGHT_GREY, line_color=RGBColor(0xBB, 0xBB, 0xBB))
add_textbox(s, "After each goal", Inches(0.5), Inches(2.6), Inches(5.4), Inches(0.4),
            font_size=15, bold=True, color=FIFA_BLUE)
add_textbox(s, "Update group standings in real time; evaluate whether next goal changes qualifiers",
            Inches(0.5), Inches(3.0), Inches(5.4), Inches(1.1),
            font_size=13, color=DARK_GREY)

add_rect(s, Inches(7.2), Inches(2.5), Inches(5.8), Inches(1.8),
         fill_color=LIGHT_GREY, line_color=RGBColor(0xBB, 0xBB, 0xBB))
add_textbox(s, "Minute by minute", Inches(7.4), Inches(2.6), Inches(5.4), Inches(0.4),
            font_size=15, bold=True, color=FIFA_BLUE)
add_textbox(s, "Track suspense on a 90-minute scale; capture how long groups remain undecided",
            Inches(7.4), Inches(3.0), Inches(5.4), Inches(1.1),
            font_size=13, color=DARK_GREY)

add_textbox(s, "⟷  Same data, two rule-sets, fully counterfactual",
            Inches(5.95), Inches(3.0), Inches(1.5), Inches(0.8),
            font_size=12, bold=True, color=UEFA_RED, align=PP_ALIGN.CENTER)

# Examples
eg_items = [
    "Euro 2024 Group F — Portugal & Turkey qualify; Czech Republic starts 0-0 vs Turkey → Czech goal flips qualification → Suspense = 1",
    "Euro 2024 Group A — Hungary scores in minute 100 vs Scotland → briefly becomes a best 3rd-place team → Suspense = 1 at that moment",
    "Two shocks evaluated: lagging team scores (+1)  OR  leading team concedes (+1)",
    "Full tie on all criteria still counts as suspenseful (next tie-break could be lots)",
]
add_bullet_box(s, eg_items,
               Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.75),
               title="Worked Examples", font_size=13.5,
               bg_color=RGBColor(0xEF, 0xF7, 0xEF))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Descriptive Statistics
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Descriptive Statistics", "Goal-level data — mean values")

# Table headers
col_positions = [Inches(2.1), Inches(3.6), Inches(5.1), Inches(6.6), Inches(8.1), Inches(9.6), Inches(11.1)]
headers = ["", "FIFA\n(WC)", "UEFA\n(WC)", "", "FIFA\n(Euro)", "UEFA\n(Euro)", ""]

add_rect(s, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.7), fill_color=FIFA_BLUE)
add_textbox(s, "Variable", Inches(0.4), Inches(1.28), Inches(1.6), Inches(0.55),
            font_size=13, bold=True, color=WHITE)
add_textbox(s, "FIFA Rules\n(World Cup)", Inches(2.9), Inches(1.28), Inches(2.0), Inches(0.55),
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "UEFA Rules\n(World Cup)", Inches(5.0), Inches(1.28), Inches(2.0), Inches(0.55),
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "FIFA Rules\n(Euro)", Inches(7.9), Inches(1.28), Inches(2.0), Inches(0.55),
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "UEFA Rules\n(Euro)", Inches(10.0), Inches(1.28), Inches(2.0), Inches(0.55),
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Qualification changed (mean)", "0.19", "0.18", "0.32", "0.29"),
    ("Qualification count (mean)",   "0.72", "0.71", "1.27", "1.11"),
    ("Suspense (mean)",              "0.44", "0.36", "0.50", "0.49"),
    ("Elo home (mean)",              "1,797", "1,797", "1,857", "1,857"),
    ("Elo away (mean)",              "1,838", "1,838", "1,896", "1,896"),
    ("Observations",                 "381",   "381",   "226",   "226"),
]
for i, (label, v1, v2, v3, v4) in enumerate(rows):
    y = Inches(2.0) + Inches(0.72 * i)
    bg = LIGHT_GREY if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.7), fill_color=bg)
    add_textbox(s, label, Inches(0.4), y + Inches(0.08), Inches(2.5), Inches(0.55),
                font_size=13.5, color=DARK_GREY,
                bold=(label == "Suspense (mean)"))
    # Highlight suspense row difference
    for col_x, val, is_fifa in [(Inches(3.0), v1, True), (Inches(5.1), v2, False),
                                  (Inches(8.0), v3, True), (Inches(10.1), v4, False)]:
        fc = DARK_GREY
        is_bold = False
        if label == "Suspense (mean)":
            fc = FIFA_BLUE if is_fifa else UEFA_RED
            is_bold = True
        add_textbox(s, val, col_x, y + Inches(0.08), Inches(1.8), Inches(0.55),
                    font_size=13.5, color=fc, align=PP_ALIGN.CENTER, bold=is_bold)

add_textbox(s,
    "WC: Suspense under FIFA (0.44) > UEFA (0.36)   |   Euro: almost identical (0.50 vs 0.49)",
    Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.55),
    font_size=14, bold=True, color=FIFA_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Statistical Tests
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Statistical Tests", "Paired t-tests and Wilcoxon signed-rank tests across all groups")

# WC results
add_rect(s, Inches(0.3), Inches(1.25), Inches(6.1), Inches(4.8),
         fill_color=RGBColor(0xE8, 0xF0, 0xFE), line_color=FIFA_BLUE, line_width=Pt(1.5))
add_textbox(s, "World Cup  (74 groups)", Inches(0.5), Inches(1.35), Inches(5.7), Inches(0.45),
            font_size=17, bold=True, color=FIFA_BLUE)

wc_rows = [
    ("Metric", "Mean Diff.", "p-value (t-test)", "p-value (Wilcoxon)"),
    ("Qualification Count", "−0.024", "0.535", "0.605"),
    ("Suspense", "−0.040", "0.008 ***", "0.010 ***"),
]
for i, row in enumerate(wc_rows):
    y = Inches(1.9) + Inches(0.72 * i)
    bg = FIFA_BLUE if i == 0 else (LIGHT_GREY if i % 2 == 1 else WHITE)
    add_rect(s, Inches(0.4), y, Inches(5.9), Inches(0.7), fill_color=bg)
    cols = [Inches(0.5), Inches(1.9), Inches(3.3), Inches(4.7)]
    for j, cell in enumerate(row):
        fc = WHITE if i == 0 else DARK_GREY
        bd = (i == 0)
        if i > 0 and j > 0 and row[0] == "Suspense":
            fc = GREEN
            bd = True
        add_textbox(s, cell, cols[j], y + Inches(0.06), Inches(1.3), Inches(0.58),
                    font_size=12, bold=bd, color=fc)

add_textbox(s, "✔ FIFA rule significantly higher suspense (p < 0.01)\n✔ No significant difference in actual qualification changes",
            Inches(0.5), Inches(3.35), Inches(5.7), Inches(1.1),
            font_size=13, color=DARK_GREY, bold=False)
add_rect(s, Inches(0.4), Inches(4.6), Inches(5.9), Inches(1.2),
         fill_color=RGBColor(0xD5, 0xE8, 0xD4), line_color=GREEN)
add_textbox(s, "FIFA raises potential qualification changes\n(suspense) without changing actual qualification\nchanges — it keeps more scenarios alive.",
            Inches(0.55), Inches(4.65), Inches(5.6), Inches(1.1),
            font_size=13, bold=True, color=RGBColor(0x1A, 0x5C, 0x1A))

# Euro results
add_rect(s, Inches(6.9), Inches(1.25), Inches(6.1), Inches(4.8),
         fill_color=RGBColor(0xFE, 0xEC, 0xE8), line_color=UEFA_RED, line_width=Pt(1.5))
add_textbox(s, "European Championship  (44 groups)", Inches(7.1), Inches(1.35), Inches(5.7), Inches(0.45),
            font_size=17, bold=True, color=UEFA_RED)

eu_rows = [
    ("Metric", "Mean Diff.", "p-value (t-test)", "p-value (Wilcoxon)"),
    ("Qualification Count", "−0.106", "0.182", "0.250"),
    ("Suspense", "−0.015", "0.389", "0.438"),
]
for i, row in enumerate(eu_rows):
    y = Inches(1.9) + Inches(0.72 * i)
    bg = UEFA_RED if i == 0 else (LIGHT_GREY if i % 2 == 1 else WHITE)
    add_rect(s, Inches(7.0), y, Inches(5.9), Inches(0.7), fill_color=bg)
    cols = [Inches(7.1), Inches(8.5), Inches(9.9), Inches(11.3)]
    for j, cell in enumerate(row):
        fc = WHITE if i == 0 else DARK_GREY
        bd = (i == 0)
        add_textbox(s, cell, cols[j], y + Inches(0.06), Inches(1.3), Inches(0.58),
                    font_size=12, bold=bd, color=fc)

add_rect(s, Inches(7.0), Inches(3.35), Inches(5.9), Inches(2.55),
         fill_color=RGBColor(0xFF, 0xEE, 0xCC), line_color=RGBColor(0xCC, 0x88, 0x00))
add_textbox(s,
    "✘ No significant difference in either metric\n\nThis serves as a placebo check:\nthe effect is specific to the World Cup,\nnot an artefact of how the rules operate\nin any tournament context.",
    Inches(7.15), Inches(3.4), Inches(5.6), Inches(2.35),
    font_size=13, color=RGBColor(0x55, 0x44, 0x00))

add_textbox(s,
    "Note: mean differences reported as UEFA − FIFA (negative = FIFA generates more suspense)",
    Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.45),
    font_size=12, italic=True, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Suspense over Time (figures)
# ═══════════════════════════════════════════════════════════════════════════════
IMG_BASE = r"c:\Users\aldi\Documents\GitHub\tb_football\paper\tb_football\images"

s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Suspense Over Time (Minute by Minute)",
             "Average suspense per minute under FIFA vs UEFA rules")

add_image(s, f"{IMG_BASE}\\susp_minute_wc.png",
          Inches(0.2), Inches(1.25), Inches(6.4), Inches(5.0))
add_image(s, f"{IMG_BASE}\\susp_minute_eu.png",
          Inches(6.7), Inches(1.25), Inches(6.4), Inches(5.0))

add_textbox(s, "World Cup", Inches(0.2), Inches(6.35), Inches(6.4), Inches(0.4),
            font_size=14, bold=True, color=FIFA_BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, "European Championship (placebo)", Inches(6.7), Inches(6.35), Inches(6.4), Inches(0.4),
            font_size=14, bold=True, color=UEFA_RED, align=PP_ALIGN.CENTER)

add_textbox(s,
    "WC: FIFA generates systematically higher suspense across all 90 minutes — gap widens near the end.  Euro: rules produce almost identical suspense profiles.",
    Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.45),
    font_size=13, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Suspense Trajectories by Group Type
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Suspense Trajectories by Group Volatility (World Cup)",
             "Stable groups = below-median qualification changes; Volatile = above median")

add_image(s, f"{IMG_BASE}\\susp_traj_wc_fifa.png",
          Inches(0.2), Inches(1.25), Inches(6.4), Inches(5.0))
add_image(s, f"{IMG_BASE}\\susp_traj_wc_uefa.png",
          Inches(6.7), Inches(1.25), Inches(6.4), Inches(5.0))

add_textbox(s, "Under FIFA rules", Inches(0.2), Inches(6.35), Inches(6.4), Inches(0.4),
            font_size=14, bold=True, color=FIFA_BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, "Under UEFA rules", Inches(6.7), Inches(6.35), Inches(6.4), Inches(0.4),
            font_size=14, bold=True, color=UEFA_RED, align=PP_ALIGN.CENTER)

add_textbox(s,
    "FIFA: stable/volatile gap is large from kick-off and persists (~0.4 pts).  UEFA: gap starts small (~0.2 pts) but widens as the match progresses.",
    Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.45),
    font_size=13, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Regression Results
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Regression Results", "Linear Probability Model — goal-level, matchday 3")

# Spec labels
add_textbox(s, "Spec 1: suspense = α + β FIFA_rule + γ t + δ t² + λ Elo_diff + year FE",
            Inches(0.3), Inches(1.25), Inches(12.7), Inches(0.4),
            font_size=13.5, italic=True, color=DARK_GREY)
add_textbox(s, "Spec 2: Spec 1 + Elo_std + (FIFA_rule × Elo_std)   |   SEs clustered at group level",
            Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.4),
            font_size=13.5, italic=True, color=DARK_GREY)

# Table
col_headers = ["Variable", "Euro Spec 1", "Euro Spec 2", "WC Spec 1", "WC Spec 2"]
col_x       = [Inches(0.35), Inches(3.1), Inches(5.25), Inches(7.4), Inches(9.55)]
col_w       = [Inches(2.7),  Inches(2.1), Inches(2.1),  Inches(2.1), Inches(2.1)]

rows_reg = [
    ("FIFA_rule",          "0.009",      "0.031",     "0.042***",  "0.075**"),
    ("",                   "(0.021)",    "(0.039)",   "(0.015)",   "(0.034)"),
    ("Minute",             "−0.007",     "−0.007",    "−0.007*",   "−0.007**"),
    ("",                   "(0.005)",    "(0.005)",   "(0.004)",   "(0.004)"),
    ("Elo diff",           "0.000",      "0.000",     "0.000",     "0.000"),
    ("Elo std",            "—",          "0.001",     "—",         "−0.002**"),
    ("FIFA × Elo std",     "—",          "−0.001",    "—",         "−0.001"),
    ("Year FE",            "Yes",        "Yes",       "Yes",       "Yes"),
    ("N",                  "452",        "450",       "762",       "760"),
]

header_y = Inches(2.2)
add_rect(s, Inches(0.3), header_y, Inches(11.5), Inches(0.55), fill_color=FIFA_BLUE)
for j, (hdr, cx, cw) in enumerate(zip(col_headers, col_x, col_w)):
    add_textbox(s, hdr, cx, header_y + Inches(0.05), cw, Inches(0.45),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows_reg):
    y = Inches(2.75) + Inches(0.42 * i)
    bg = LIGHT_GREY if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.3), y, Inches(11.5), Inches(0.41), fill_color=bg)
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        fc = DARK_GREY
        bd = False
        if j in (3, 4) and "***" in cell:
            fc = GREEN; bd = True
        elif j in (3, 4) and "**" in cell and "***" not in cell:
            fc = RGBColor(0x00, 0x70, 0x00); bd = True
        add_textbox(s, cell, cx, y + Inches(0.02), cw, Inches(0.37),
                    font_size=12.5, color=fc, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                    bold=bd)

add_textbox(s,
    "WC: +4.2 pp (Spec 1) / +7.5 pp (Spec 2) — effect is structural, independent of competitive balance (FIFA×Elo std insignificant)\n"
    "Euro: insignificant in both specs — consistent with placebo interpretation",
    Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.65),
    font_size=13, bold=True, color=FIFA_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Key Findings Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Key Findings")

findings = [
    ("1", "FIFA's goal-difference rule generates significantly more suspense in the World Cup",
     "A goal on matchday 3 under FIFA rules is 4.2 pp more likely to occur in a suspenseful state (p < 0.01); robust to competitive balance controls.",
     GREEN),
    ("2", "The effect is structural — not driven by competitive imbalance",
     "Interaction FIFA × Elo std is statistically insignificant: FIFA advantage is uniform across balanced and unbalanced groups.",
     FIFA_BLUE),
    ("3", "Suspense under FIFA rules is more stable and persistent over time",
     "FIFA produces higher suspense throughout the full 90 minutes; the gap is largest near the end. Volatile groups are more suspenseful under both rules.",
     ACCENT),
    ("4", "European Championship as a placebo: no significant effect",
     "Neither t-test nor Wilcoxon test detects a difference in Euro groups (p > 0.38). Effect is specific to the World Cup's goals distribution.",
     UEFA_RED),
]

for i, (num, title, detail, color) in enumerate(findings):
    y = Inches(1.25) + Inches(1.35 * i)
    add_rect(s, Inches(0.3), y, Inches(0.65), Inches(1.1), fill_color=color)
    add_textbox(s, num, Inches(0.35), y + Inches(0.25), Inches(0.55), Inches(0.6),
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.95), y, Inches(12.05), Inches(1.1),
             fill_color=LIGHT_GREY, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(s, title, Inches(1.1), y + Inches(0.05), Inches(11.7), Inches(0.45),
                font_size=15, bold=True, color=color)
    add_textbox(s, detail, Inches(1.1), y + Inches(0.5), Inches(11.7), Inches(0.55),
                font_size=13, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Conclusion & Policy Implications
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
slide_header(s, "Conclusion & Policy Implications")

# Main take-away
add_rect(s, Inches(0.3), Inches(1.25), Inches(12.7), Inches(1.35), fill_color=FIFA_BLUE)
add_textbox(s,
    "Institutional design matters: the choice of tie-breaking criterion has measurable consequences\n"
    "for the suspense experienced by billions of fans who follow the World Cup.",
    Inches(0.5), Inches(1.33), Inches(12.3), Inches(1.2),
    font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

conclusion_items = [
    "Goal-difference (FIFA) keeps more qualification scenarios alive simultaneously → more persistent suspense",
    "Head-to-head (UEFA) locks in outcomes early once previous fixtures are played → lower potential suspense",
    "Audiences aware of qualification stakes respond to this uncertainty → direct link to entertainment value",
    "2026 FIFA World Cup will switch to head-to-head as primary tiebreaker → our results predict a decline in WC group-stage suspense",
    "Governing bodies should weigh suspension-generating properties of rules alongside fairness considerations",
]
add_bullet_box(s, conclusion_items,
               Inches(0.3), Inches(2.75), Inches(12.7), Inches(3.55),
               title="Take-aways for Governing Bodies", font_size=14.5,
               bg_color=LIGHT_GREY)

add_rect(s, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.75),
         fill_color=RGBColor(0xFF, 0xF0, 0xC0), line_color=RGBColor(0xCC, 0x88, 0x00))
add_textbox(s,
    "Paper status: all sections complete (Intro · Lit Review · Data · Methodology · Results · Conclusion · Appendix)\n"
    "Next steps: final proofread, journal submission target TBD",
    Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.65),
    font_size=13, color=RGBColor(0x55, 0x44, 0x00))


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\aldi\Documents\GitHub\tb_football\presentation_supervisor.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
